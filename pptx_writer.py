"""
pptx_writer.py  —  Populate slides 3, 4, 5 of the POE template.

Golden rule: NEVER remove or overwrite any shape whose top edge sits
             within the top HEADER_ZONE_IN inches of the slide.

Slide 3  →  URL text box  +  screenshot image  (below header)
Slide 4  →  Results table image fills entire body area  (below header)
Slide 5  →  2×2 grid of 4 LinkedIn/Advert images  (below header)
"""

import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Shapes whose top edge is above this line belong to the header — never touch them
HEADER_ZONE_IN = 1.4


# ─────────────────────────────────────────────────────────────────────────────
# Low-level helpers
# ─────────────────────────────────────────────────────────────────────────────

def _in(emu: int) -> float:
    """EMU → inches."""
    return emu / 914400.0


def _is_header(shape) -> bool:
    """True if this shape lives inside the header band."""
    return _in(shape.top) < HEADER_ZONE_IN


def _shape_text(shape) -> str:
    if shape.has_text_frame:
        return " ".join(
            r.text for p in shape.text_frame.paragraphs for r in p.runs
        ).strip()
    return ""


def _remove(shape):
    sp = shape._element
    sp.getparent().remove(sp)


def _add_image(slide, path: str, left: float, top: float, w: float, h: float):
    """Place an image at (left, top) with size (w, h) all in inches."""
    if path and os.path.isfile(path):
        slide.shapes.add_picture(
            path, Inches(left), Inches(top), Inches(w), Inches(h)
        )


def _is_placeholder_shape(shape) -> bool:
    """
    Returns True only for shapes that are clearly placeholder/empty body content
    that should be replaced.  Header shapes are excluded by the caller.
    """
    txt = _shape_text(shape)
    # Empty body shape
    if not txt:
        return True
    # Explicit placeholder token patterns
    if re.search(r'\{\{|\[\[|placeholder|_here\b', txt, re.I):
        return True
    return False


def _clear_body_placeholders(slide):
    """
    Remove non-header placeholder shapes from the slide body.
    Leaves all header shapes, text shapes with real content, and decorative shapes intact.
    """
    to_remove = [
        s for s in slide.shapes
        if not _is_header(s) and _is_placeholder_shape(s)
    ]
    for s in to_remove:
        try:
            _remove(s)
        except Exception:
            pass


# ─────────────────────────────────────────────────────────────────────────────
# Slide 3 — URL + Screenshot
# ─────────────────────────────────────────────────────────────────────────────

def _slide3(slide, prs, extraction: dict, analysis: dict):
    sw = _in(prs.slide_width)
    sh = _in(prs.slide_height)

    url  = extraction.get("slide3_url") or (analysis.get("urls") or [""])[0]
    img3 = extraction.get("slide3_screenshot")

    _clear_body_placeholders(slide)

    # ── URL text box just below header ────────────────────────────────────────
    URL_TOP   = HEADER_ZONE_IN + 0.15
    URL_H     = 0.4
    MARGIN    = 0.4

    if url:
        txb = slide.shapes.add_textbox(
            Inches(MARGIN), Inches(URL_TOP),
            Inches(sw - MARGIN * 2), Inches(URL_H)
        )
        tf = txb.text_frame
        tf.word_wrap = False
        para = tf.paragraphs[0]
        run  = para.add_run()
        run.text = url
        run.font.size      = Pt(12)
        run.font.bold      = False
        run.font.underline = True
        run.font.color.rgb = RGBColor(0x00, 0x56, 0xB3)

    # ── Screenshot fills rest of body ─────────────────────────────────────────
    if img3:
        IMG_TOP = URL_TOP + URL_H + 0.12
        IMG_H   = sh - IMG_TOP - 0.25
        IMG_W   = sw - MARGIN * 2
        _add_image(slide, img3, MARGIN, IMG_TOP, IMG_W, IMG_H)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 4 — Results table image
# ─────────────────────────────────────────────────────────────────────────────

def _slide4(slide, prs, extraction: dict):
    sw = _in(prs.slide_width)
    sh = _in(prs.slide_height)

    table_img = extraction.get("slide4_table_img")
    if not table_img:
        return

    _clear_body_placeholders(slide)

    MARGIN  = 0.35
    TOP     = HEADER_ZONE_IN + 0.1
    IMG_W   = sw - MARGIN * 2
    IMG_H   = sh - TOP - 0.25

    _add_image(slide, table_img, MARGIN, TOP, IMG_W, IMG_H)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 5 — 2×2 advert image grid
# ─────────────────────────────────────────────────────────────────────────────

def _slide5(slide, prs, advert_images: list[str]):
    sw = _in(prs.slide_width)
    sh = _in(prs.slide_height)

    if not advert_images:
        return

    _clear_body_placeholders(slide)

    MARGIN   = 0.35
    GAP      = 0.18
    COLS     = 2
    ROWS     = 2
    TOP      = HEADER_ZONE_IN + 0.1
    BODY_W   = sw - MARGIN * 2
    BODY_H   = sh - TOP - 0.25

    cell_w = (BODY_W - GAP * (COLS - 1)) / COLS
    cell_h = (BODY_H - GAP * (ROWS - 1)) / ROWS

    for idx, img_path in enumerate(advert_images[:4]):
        col  = idx % COLS
        row  = idx // COLS
        left = MARGIN + col * (cell_w + GAP)
        top  = TOP    + row * (cell_h + GAP)
        _add_image(slide, img_path, left, top, cell_w, cell_h)


# ─────────────────────────────────────────────────────────────────────────────
# Entry point
# ─────────────────────────────────────────────────────────────────────────────

def populate_pptx(
    template_path: str,
    output_path:   str,
    analysis:      dict,
    extraction:    dict,
) -> None:

    prs    = Presentation(template_path)
    slides = prs.slides

    if len(slides) < 5:
        raise ValueError(
            f"Template has only {len(slides)} slide(s) — "
            "slides 3, 4 and 5 are required."
        )

    _slide3(slides[2], prs, extraction, analysis)
    _slide4(slides[3], prs, extraction)
    _slide5(slides[4], prs, extraction.get("advert_images", []))

    prs.save(output_path)
